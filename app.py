import os
import traceback
from flask import Flask, render_template, request, redirect, send_from_directory, jsonify, after_this_request
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.util import Pt
from werkzeug.utils import secure_filename
import gc
import time
from deep_translator import GoogleTranslator
import signal
from functools import wraps
from threading import Event
import tempfile
import shutil
import threading
import psutil
import sys
from pympler import summary, muppy

UPLOAD_FOLDER = 'uploads'
CONVERTED_FOLDER = 'converted'
CHUNK_FOLDER = 'chunks'
MAX_FILE_AGE = 300  # 5 minutes in seconds
CHUNK_SIZE = 512 * 1024  # 512KB chunk size for uploads - increased for better performance
MAX_SLIDES_PER_BATCH = 2  # Process 2 slides at a time
MEMORY_CLEANUP_DELAY = 0.5  # 0.5 second delay between memory cleanups - reduced for better performance

# Create necessary directories
for folder in [UPLOAD_FOLDER, CONVERTED_FOLDER, CHUNK_FOLDER]:
    os.makedirs(folder, exist_ok=True)

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['CONVERTED_FOLDER'] = CONVERTED_FOLDER
app.config['CHUNK_FOLDER'] = CHUNK_FOLDER
app.config['MAX_CONTENT_LENGTH'] = None  # Remove global limit
app.config['MAX_CHUNK_SIZE'] = CHUNK_SIZE
app.config['SEND_FILE_MAX_AGE_DEFAULT'] = 0  # Disable caching
app.config['TEMPLATES_AUTO_RELOAD'] = False  # Disable template auto-reload
app.config['JSON_AS_ASCII'] = False  # Support non-ASCII characters
app.config['PERMANENT_SESSION_LIFETIME'] = 1800  # 30 minutes
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024  # 100MB max file size
app.config['ALLOWED_EXTENSIONS'] = {'pptx'}  # Update allowed extensions
app.config['REQUEST_TIMEOUT'] = 600  # 10 minutes timeout
app.config['THREADED'] = True

# Add worker configuration
worker_class = 'gthread'
worker_connections = 1000
worker_timeout = 600  # 10 minutes
threads = 4
max_requests = 1000
max_requests_jitter = 50

# Global abort event
abort_event = Event()

def check_abort():
    """Check if the process should be aborted"""
    return abort_event.is_set()

def reset_abort():
    """Reset the abort event"""
    abort_event.clear()

def handle_abort(signal, frame):
    """Signal handler for abort"""
    abort_event.set()

# Register signal handler
signal.signal(signal.SIGINT, handle_abort)

def log_memory_usage(tag=""):
    """Log current memory usage"""
    process = psutil.Process(os.getpid())
    memory_info = process.memory_info()
    print(f"[Memory {tag}] RSS: {memory_info.rss / (1024 * 1024):.2f} MB, VMS: {memory_info.vms / (1024 * 1024):.2f} MB")

def force_memory_cleanup():
    """Force memory cleanup with delay"""
    if check_abort():
        raise Exception("Process aborted by user")
    
    log_memory_usage("Before GC")
    gc.collect()
    log_memory_usage("After GC")
    time.sleep(MEMORY_CLEANUP_DELAY)

def cleanup_old_files():
    """Clean up files older than MAX_FILE_AGE seconds"""
    current_time = time.time()
    
    for folder in [UPLOAD_FOLDER, CONVERTED_FOLDER, CHUNK_FOLDER]:
        for filename in os.listdir(folder):
            filepath = os.path.join(folder, filename)
            if os.path.isfile(filepath):
                file_age = current_time - os.path.getmtime(filepath)
                if file_age > MAX_FILE_AGE:
                    try:
                        os.remove(filepath)
                        print(f"[Cleanup] Removed old file: {filepath}")
                    except Exception as e:
                        print(f"[Cleanup] Error removing {filepath}: {e}")

def log_error(error, context=""):
    """Helper function to log errors with context"""
    print(f"[Error] {context}: {str(error)}")
    print(f"[Error] Stack trace: {traceback.format_exc()}")

def convert_number_to_arabic(text):
    """Helper function to convert numbers to Arabic and clean up formatting"""
    text = text.strip()
    if text.endswith('.'):
        text = text[:-1]
    if text.startswith('.'):
        text = text[1:]
        
    arabic_text = ""
    last_was_digit = False
    
    for char in text:
        if char.isdigit():
            arabic_text += chr(ord('٠') + int(char))
            last_was_digit = True
        else:
            if char == '.' and last_was_digit:
                continue
            arabic_text += char
            last_was_digit = False
            
    return arabic_text.strip()

def process_text_frame_format(text_frame, direction):
    """Process text frame formatting only"""
    try:
        is_placeholder = hasattr(text_frame, 'is_placeholder') and text_frame.is_placeholder
        
        if direction == 'en_to_ar':
            # Get or create the text body element
            if hasattr(text_frame, '_element'):
                txBody = text_frame._element
            else:
                txBody = text_frame
            
            # Set RTL for the text body
            if not hasattr(txBody, 'bodyPr'):
                txBody.insert(0, parse_xml(r'<p:bodyPr xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>'))
            txBody.bodyPr.set('rtl', '1')
            
            for paragraph in text_frame.paragraphs:
                # Set paragraph properties
                if not is_placeholder:
                    paragraph.alignment = PP_ALIGN.RIGHT
                    
                    # Handle numbering
                    if hasattr(paragraph._pPr, 'numPr') and paragraph._pPr.numPr is not None:
                        pPr = paragraph._element.get_or_add_pPr()
                        if pPr.numPr is not None:
                            pPr.set('rtl', '1')
                            lvl = pPr.get_or_add_numPr().get_or_add_ilvl()
                            lvl.val = 0
                    
                    # Set RTL for paragraph
                    pPr = paragraph._element.get_or_add_pPr()
                    pPr.set('rtl', '1')
                    
                    # Add RTL element
                    rtl = parse_xml(r'<a:rtl xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">1</a:rtl>')
                    paragraph._p.insert(0, rtl)
                
                # Process each run in the paragraph
                for run in paragraph.runs:
                    # Set font properties
                    run.font.name = "Traditional Arabic"
                    if not run.font.size:
                        run.font.size = Pt(18)
                    
                    # Convert numbers to Arabic
                    if any(char.isdigit() for char in run.text):
                        run.text = convert_number_to_arabic(run.text)
                    
                    # Add Unicode RTL markers to the text
                    text = run.text
                    if text and not text.startswith('\u200F'):  # RTL Mark
                        text = '\u200F' + text
                    if text and not text.endswith('\u200F'):  # RTL Mark
                        text = text + '\u200F'
                    run.text = text
                    
                    # Set language and RTL properties for run
                    rPr = run._r.get_or_add_rPr()
                    rPr.set(qn('w:lang'), 'ar-SA')
                    
        else:  # ar_to_en
            # Get or create the text body element
            if hasattr(text_frame, '_element'):
                txBody = text_frame._element
            else:
                txBody = text_frame
            
            # Set LTR for the text body
            if not hasattr(txBody, 'bodyPr'):
                txBody.insert(0, parse_xml(r'<p:bodyPr xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>'))
            txBody.bodyPr.set('rtl', '0')
            
            for paragraph in text_frame.paragraphs:
                # Set paragraph properties
                if not is_placeholder:
                    paragraph.alignment = PP_ALIGN.LEFT
                    
                    # Handle numbering
                    if hasattr(paragraph._pPr, 'numPr') and paragraph._pPr.numPr is not None:
                        pPr = paragraph._element.get_or_add_pPr()
                        if pPr.numPr is not None:
                            pPr.set('rtl', '0')
                    
                    # Set LTR for paragraph
                    pPr = paragraph._element.get_or_add_pPr()
                    pPr.set('rtl', '0')
                    
                    # Remove RTL elements
                    rtl_elements = paragraph._p.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}rtl')
                    for rtl_elem in rtl_elements:
                        paragraph._p.remove(rtl_elem)
                
                # Process each run in the paragraph
                for run in paragraph.runs:
                    # Set font properties
                    run.font.name = "Arial"
                    if not run.font.size:
                        run.font.size = Pt(12)
                    
                    # Remove Unicode RTL markers from the text
                    text = run.text
                    text = text.replace('\u200F', '')  # Remove RTL Mark
                    run.text = text
                    
                    # Set language and LTR properties for run
                    rPr = run._r.get_or_add_rPr()
                    rPr.set(qn('w:lang'), 'en-US')
                    
    except Exception as e:
        log_error(e, "Error processing text frame format")

def process_shape_format(shape, slide_width, direction, in_group=False):
    """Process shape formatting only"""
    try:
        is_placeholder = hasattr(shape, 'is_placeholder') and shape.is_placeholder
        
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            if not in_group and not is_placeholder:
                try:
                    shape.left = slide_width - shape.left - shape.width
                except Exception as e:
                    log_error(e, "Error mirroring group container")
            for child in shape.shapes:
                process_shape_format(child, slide_width, direction, in_group=True)
        else:
            if not in_group and not is_placeholder:
                try:
                    shape.left = slide_width - shape.left - shape.width
                except Exception as e:
                    log_error(e, "Error mirroring shape")
            if shape.has_text_frame:
                process_text_frame_format(shape.text_frame, direction)
    except Exception as e:
        log_error(e, "Error processing shape format")

def convert_pptx(input_path, output_path, slide_indices=None, direction='en_to_ar'):
    try:
        print(f"[Conversion] Starting conversion from {input_path} to {output_path}")
        
        # Reset abort flag at start of conversion
        reset_abort()
        
        # Initialize translator
        translator = GoogleTranslator(source='en', target='ar') if direction == 'en_to_ar' else GoogleTranslator(source='ar', target='en')
        
        # Load presentation with minimal memory usage
        prs = Presentation(input_path)
        slide_width = prs.slide_width
        total_slides = len(prs.slides)
        
        print(f"[Conversion] Total slides: {total_slides}")
        
        if slide_indices:
            slide_indices = [i for i in slide_indices if 1 <= i <= total_slides]
            print(f"[Conversion] Processing specific slides: {slide_indices}")
        else:
            print(f"[Conversion] Processing all slides")
        
        # Translation cache to avoid repeating translations
        translation_cache = {}
        
        # Process slides in smaller batches to reduce memory usage
        batch_size = 3  # Process 3 slides at a time
        
        # Process all slides for formatting and translation
        for batch_start in range(1, total_slides + 1, batch_size):
            batch_end = min(batch_start + batch_size - 1, total_slides)
            print(f"[Conversion] Processing slide batch {batch_start}-{batch_end}")
            
            # Save progress after each batch
            should_save = False
            
            for slide_index in range(batch_start, batch_end + 1):
                # Check for abort signal
                if check_abort():
                    print("[Conversion] Process aborted by user")
                    raise Exception("Process aborted by user")
                    
                if slide_indices and slide_index not in slide_indices:
                    continue
                    
                print(f"[Conversion] Processing slide {slide_index}/{total_slides}")
                slide = prs.slides[slide_index - 1]  # 0-based index
                
                # Process each shape in the slide
                for shape in slide.shapes:
                    if check_abort():
                        print("[Conversion] Process aborted by user")
                        raise Exception("Process aborted by user")
                    
                    # Process formatting
                    process_shape_format(shape, slide_width, direction)
                    
                    # Process translation if shape has text
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text = run.text.strip()
                                if text and text is not None:  # Only translate non-empty, non-None text
                                    try:
                                        # Check cache first
                                        if text in translation_cache:
                                            run.text = translation_cache[text]
                                        else:
                                            translated_text = translator.translate(text)
                                            if translated_text and translated_text is not None:
                                                run.text = translated_text
                                                # Cache the translation
                                                translation_cache[text] = translated_text
                                            else:
                                                print(f"[Translation Warning] Empty translation result for text: {text}")
                                    except Exception as e:
                                        print(f"[Translation Error] Failed to translate text: {e}")
                                        continue
                
                # Mark that we should save progress
                should_save = True
            
            # Force memory cleanup after each batch
            force_memory_cleanup()
            
            # Save progress after each batch
            if should_save:
                try:
                    prs.save(output_path)
                    print(f"[Conversion] Saved progress after slide batch {batch_start}-{batch_end}")
                except Exception as e:
                    log_error(e, "Error saving progress")
            
            # Clear some memory between batches
            del translation_cache
            translation_cache = {}
            gc.collect()
        
        return 'completed'
    except Exception as e:
        if str(e) == "Process aborted by user":
            return 'aborted'
        log_error(e, "Error during PowerPoint conversion")
        raise
    finally:
        # Ensure we clean up the input file and force garbage collection
        try:
            if os.path.exists(input_path):
                os.remove(input_path)
            force_memory_cleanup()
        except Exception as e:
            log_error(e, "Error cleaning up input file")

def assemble_chunks(chunk_files, output_path):
    """Assemble uploaded chunks into a single file"""
    try:
        with open(output_path, 'wb') as outfile:
            for chunk_file in sorted(chunk_files):
                with open(chunk_file, 'rb') as infile:
                    outfile.write(infile.read())
                os.remove(chunk_file)  # Delete chunk after use
        return True
    except Exception as e:
        log_error(e, "Error assembling chunks")
        return False

@app.before_request
def check_request_size():
    """Check request size before processing"""
    if request.path == '/upload-chunk':
        return  # Skip size check for chunk uploads

def cleanup_on_startup():
    """Clean up all temporary files on startup"""
    print("[Startup] Cleaning up temporary directories...")
    try:
        # Clean up upload folder
        for folder in [UPLOAD_FOLDER, CONVERTED_FOLDER, CHUNK_FOLDER]:
            if os.path.exists(folder):
                for filename in os.listdir(folder):
                    filepath = os.path.join(folder, filename)
                    try:
                        if os.path.isfile(filepath):
                            os.remove(filepath)
                            print(f"[Startup] Removed file: {filepath}")
                        elif os.path.isdir(filepath):
                            shutil.rmtree(filepath)
                            print(f"[Startup] Removed directory: {filepath}")
                    except Exception as e:
                        print(f"[Startup] Error removing {filepath}: {e}")
        print("[Startup] Cleanup complete")
    except Exception as e:
        print(f"[Startup] Error during cleanup: {e}")

# Run cleanup at startup
cleanup_on_startup()

@app.route('/', methods=['GET'])
def index():
    cleanup_old_files()  # Clean up old files on page load
    return render_template('index.html')

def allowed_file(filename):
    """Check if the file extension is allowed"""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

@app.route('/upload-chunk', methods=['POST'])
def upload_chunk():
    try:
        if 'file' not in request.files:
            print("[Upload] No file part in request")
            return jsonify({'error': 'No file part'}), 400
        
        file = request.files['file']
        chunk_index = int(request.form.get('chunk_index', 0))
        total_chunks = int(request.form.get('total_chunks', 1))
        
        # Get original filename from form data
        original_filename = request.form.get('original_filename', '')
        if not original_filename:
            print("[Upload] No original filename provided")
            return jsonify({'error': 'No original filename provided'}), 400
            
        # Sanitize filename consistently
        sanitized_filename = secure_filename(original_filename)
        print(f"[Upload] Original filename: {original_filename}, Sanitized: {sanitized_filename}")

        if not allowed_file(sanitized_filename):
            print("[Upload] Invalid file type:", sanitized_filename)
            return jsonify({'error': 'File type not allowed'}), 400

        print(f"[Upload] Processing chunk {chunk_index + 1}/{total_chunks} for file: {sanitized_filename}")

        # Create a temporary directory for chunks if it doesn't exist
        chunk_dir = os.path.join(app.config['CHUNK_FOLDER'], sanitized_filename)
        os.makedirs(chunk_dir, exist_ok=True)

        # Save the chunk
        chunk_path = os.path.join(chunk_dir, f'chunk_{chunk_index}')
        file.save(chunk_path)
        print(f"[Upload] Saved chunk {chunk_index} to: {chunk_path}")

        # If this is the last chunk, combine all chunks
        if chunk_index == total_chunks - 1:
            print("[Upload] Last chunk received, combining chunks...")
            final_path = os.path.join(app.config['UPLOAD_FOLDER'], sanitized_filename)
            
            try:
                with open(final_path, 'wb') as outfile:
                    for i in range(total_chunks):
                        chunk_path = os.path.join(chunk_dir, f'chunk_{i}')
                        with open(chunk_path, 'rb') as infile:
                            outfile.write(infile.read())
                print(f"[Upload] Successfully combined chunks into: {final_path}")

                # Clean up chunks
                shutil.rmtree(chunk_dir)
                print("[Upload] Cleaned up chunk directory")

                return jsonify({
                    'message': 'File upload complete',
                    'filename': sanitized_filename
                }), 200
            except Exception as e:
                print(f"[Upload] Error combining chunks: {str(e)}")
                return jsonify({'error': f'Error combining chunks: {str(e)}'}), 500

        return jsonify({
            'message': f'Chunk {chunk_index + 1}/{total_chunks} uploaded successfully'
        }), 200

    except Exception as e:
        print(f'[Error] Chunk upload failed: {str(e)}')
        return jsonify({'error': str(e)}), 500

@app.route('/abort', methods=['POST'])
def abort_conversion():
    """Endpoint to abort the conversion process"""
    try:
        abort_event.set()
        return jsonify({'status': 'success', 'message': 'Abort signal sent'})
    except Exception as e:
        return jsonify({'status': 'error', 'message': str(e)}), 500

@app.route('/convert', methods=['POST'])
def convert():
    input_path = None
    output_path = None
    
    try:
        print("[Convert] Starting conversion process")
        log_memory_usage("Start Conversion")
        
        # Reset abort flag at start of conversion
        reset_abort()
        
        # Get the original filename from form data
        original_filename = request.form.get('original_filename')
        if not original_filename:
            print("[Convert] No original filename provided")
            return jsonify({'status': 'error', 'message': 'No original filename provided'}), 400
        
        # Important: Use secure_filename to ensure the same filename format used during upload
        sanitized_filename = secure_filename(original_filename)
        print(f"[Convert] Original filename: {original_filename}, Sanitized: {sanitized_filename}")
            
        # Use the assembled file from the upload folder with sanitized filename
        input_path = os.path.join(app.config['UPLOAD_FOLDER'], sanitized_filename)
        if not os.path.exists(input_path):
            print("[Convert] Input file not found:", input_path)
            # Try to find the file with different filename variations
            possible_files = os.listdir(app.config['UPLOAD_FOLDER'])
            print(f"[Convert] Available files in upload folder: {possible_files}")
            
            # Check if there's a similarly named file
            similar_files = [f for f in possible_files if f.lower().startswith(sanitized_filename.lower().split('.')[0])]
            if similar_files:
                print(f"[Convert] Found similar files: {similar_files}")
                input_path = os.path.join(app.config['UPLOAD_FOLDER'], similar_files[0])
                print(f"[Convert] Using alternative file: {input_path}")
            else:
                return jsonify({'status': 'error', 'message': 'Input file not found'}), 400

        output_name = request.form.get('outputName')
        if not output_name:
            print("[Convert] No output name provided")
            return jsonify({'status': 'error', 'message': 'No output name provided'}), 400

        print("[Convert] Processing file:", os.path.basename(input_path))

        slide_nums_raw = request.form.get('slideNumbers', '')
        conversion_direction = request.form.get('conversionDirection', 'en_to_ar')
        enable_translation = request.form.get('translationToggle', 'true').lower() == 'true'
        print("[Convert] Conversion direction:", conversion_direction)
        print("[Convert] Translation enabled:", enable_translation)

        # Process slide numbers
        slide_indices = None
        if slide_nums_raw.strip():
            try:
                slide_indices = [int(num.strip()) for num in slide_nums_raw.split(',') if num.strip().isdigit()]
                print("[Convert] Processing slides:", slide_indices)
            except ValueError:
                print("[Convert] Invalid slide numbers format")
                slide_indices = None

        output_filename = secure_filename(output_name) + '.pptx'
        output_path = os.path.join(app.config['CONVERTED_FOLDER'], output_filename)
        print("[Convert] Output path:", output_path)

        # Convert with memory optimization
        log_memory_usage("Before Conversion")
        gc.collect()  # Force garbage collection before processing
        print("[Convert] Starting conversion process")
        
        # Run the conversion in a separate thread to prevent worker timeout
        def run_conversion():
            nonlocal input_path, output_path, slide_indices, conversion_direction
            try:
                status = convert_pptx(
                    input_path=input_path,
                    output_path=output_path,
                    slide_indices=slide_indices,
                    direction=conversion_direction
                )
                print("[Convert] Conversion status:", status)
                log_memory_usage("After Conversion")
                
                # Clean up input file after successful conversion
                try:
                    if os.path.exists(input_path):
                        os.remove(input_path)
                        print("[Convert] Input file cleaned up")
                    gc.collect()
                except Exception as e:
                    print("[Convert] Error cleaning up input file:", str(e))
                    log_error(e, "Error cleaning up input file")
            except Exception as e:
                print("[Convert] Error in conversion thread:", str(e))
                log_error(e, "Error in conversion thread")
        
        conversion_thread = threading.Thread(target=run_conversion)
        conversion_thread.daemon = True
        conversion_thread.start()
        
        # Wait for the thread to finish with a timeout
        conversion_thread.join(timeout=600)  # 10 minute timeout
        
        if conversion_thread.is_alive():
            # Thread is still running, which means it's taking too long
            abort_event.set()  # Signal thread to abort
            return jsonify({
                'status': 'error',
                'message': 'Conversion is taking too long and has been aborted'
            }), 500

        print("[Convert] Conversion completed successfully")
        return jsonify({
            'status': 'completed',
            'download_url': f'/download/{output_filename}'
        })

    except Exception as e:
        print("[Convert] Error during conversion:", str(e))
        # Clean up files in case of error
        try:
            if input_path and os.path.exists(input_path):
                os.remove(input_path)
            if output_path and os.path.exists(output_path):
                os.remove(output_path)
            gc.collect()  # Force garbage collection after error
            log_memory_usage("After Error")
        except Exception as cleanup_error:
            print("[Convert] Error during cleanup:", str(cleanup_error))
            log_error(cleanup_error, "Error cleaning up files after error")
            
        log_error(e, "Error during request processing")
        return jsonify({'status': 'error', 'message': str(e)}), 500

def delayed_delete(file_path, delay=10):
    """Delete a file after a delay to ensure it's no longer in use"""
    def delete_file():
        time.sleep(delay)  # Wait for download to complete
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
                print(f"[Cleanup] Successfully deleted file: {file_path}")
        except Exception as e:
            print(f"[Cleanup Error] Failed to delete file {file_path}: {e}")
    
    # Start deletion in a separate thread
    thread = threading.Thread(target=delete_file)
    thread.daemon = True  # Thread will be terminated when main program exits
    thread.start()

@app.route('/download/<filename>')
def download_file(filename):
    try:
        file_path = os.path.join(app.config['CONVERTED_FOLDER'], filename)
        
        if not os.path.exists(file_path):
            return jsonify({'status': 'error', 'message': 'File not found'}), 404

        # Send the file
        response = send_from_directory(
            app.config['CONVERTED_FOLDER'],
            filename,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

        # Schedule file deletion after download
        delayed_delete(file_path)

        return response

    except Exception as e:
        log_error(e, "Error during file download")
        return jsonify({'status': 'error', 'message': str(e)}), 500

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5005))
    app.run(host='0.0.0.0', port=port, debug=False, threaded=True)
